#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Bekaert 1+6+12 BFM4-4 Wire Process 项目 PPT 生成器
按 DP 阶段生成汇报模板
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_title_slide(prs, title, subtitle, date="2026-07-15"):
    """创建封面页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)

    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 0, 0)
    title_para.alignment = PP_ALIGN.CENTER

    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = RGBColor(60, 60, 60)
    subtitle_para.alignment = PP_ALIGN.CENTER

    # 日期
    date_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(0.5))
    date_frame = date_box.text_frame
    date_frame.text = f"汇报日期：{date}"
    date_para = date_frame.paragraphs[0]
    date_para.font.size = Pt(18)
    date_para.font.color.rgb = RGBColor(120, 120, 120)
    date_para.alignment = PP_ALIGN.CENTER

    return slide

def create_content_slide(prs, title, content_items, highlight=None):
    """创建内容页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 0, 0)

    # 内容框
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    for i, item in enumerate(content_items):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()

        p.text = item['text']
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(60, 60, 60)

        if item.get('bold', False):
            p.font.bold = True

        if item.get('indent', False):
            p.level = 1

        if item.get('highlight', False):
            p.font.color.rgb = RGBColor(200, 50, 50)

    return slide

def create_table_slide(prs, title, headers, rows):
    """创建表格页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 0, 0)

    # 表格
    table = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(0.5), Inches(1.5), Inches(9), Inches(5)).table
    # 设置表头背景
    for i in range(len(headers)):
        table.cell(0, i).fill.solid()
        table.cell(0, i).fill.fore_color.rgb = RGBColor(230, 230, 230)
        table.cell(0, i).text_frame.paragraphs[0].font.bold = True
        table.cell(0, i).text_frame.paragraphs[0].font.size = Pt(18)
        table.cell(0, i).text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        table.cell(0, i).alignment = PP_ALIGN.CENTER

    # 填充表头内容
    for i, header in enumerate(headers):
        table.cell(0, i).text = header

    # 填充数据行
    for i, row in enumerate(rows):
        for j, cell_value in enumerate(row):
            table.cell(i + 1, j).text = cell_value
            table.cell(i + 1, j).text_frame.paragraphs[0].font.size = Pt(18)
            table.cell(i + 1, j).text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            table.cell(i + 1, j).alignment = PP_ALIGN.CENTER

    return slide

def create_summary_slide(prs, title, summary_items, recommendations=None):
    """创建总结页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 0, 0)

    # 左侧：总结
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.5), Inches(5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True

    left_frame.text = "📊 项目总结"
    left_frame.paragraphs[0].font.bold = True
    left_frame.paragraphs[0].font.size = Pt(24)
    left_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    for i, item in enumerate(summary_items):
        if i == 0:
            p = left_frame.paragraphs[0]
        else:
            p = left_frame.add_paragraph()

        p.text = f"• {item}"
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(60, 60, 60)

    # 右侧：建议
    right_box = slide.shapes.add_textbox(Inches(6.5), Inches(1.5), Inches(3), Inches(5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True

    right_frame.text = "💡 建议"
    right_frame.paragraphs[0].font.bold = True
    right_frame.paragraphs[0].font.size = Pt(24)
    right_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    for i, item in enumerate(recommendations):
        if i == 0:
            p = right_frame.paragraphs[0]
        else:
            p = right_frame.add_paragraph()

        p.text = f"• {item}"
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(60, 60, 60)

    return slide

def create_dp1_slide(prs):
    """DP1: 项目启动"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "DP1: 项目启动"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 0, 0)

    # 内容
    content_items = [
        {"text": "项目立项", "bold": True},
        {"text": "• 明确项目目标与范围"},
        {"text": "• 组建项目团队"},
        {"text": "• 确定关键里程碑"},
        {"text": "• 获得管理层批准"},
        {"text": ""},
        {"text": "📅 截止日期：2025-06-30", "bold": True, "highlight": True},
        {"text": "✅ 状态：已完成"},
    ]

    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    for i, item in enumerate(content_items):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()

        p.text = item['text']
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(60, 60, 60)
        p.font.bold = item.get('bold', False)

    return slide

def create_dp3_slide(prs):
    """DP3: 概念冻结"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "DP3: 概念冻结"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 0, 0)

    content_items = [
        {"text": "技术方案确认", "bold": True},
        {"text": "• BFM4-4 线材工艺方案"},
        {"text": "• 设备配置确认（4台BFM4 + 2台AWC 2.0 WND）"},
        {"text": "• 产能规划（175吨）"},
        {"text": ""},
        {"text": "市场验证", "bold": True},
        {"text": "• 客户需求确认"},
        {"text": "• 商业案例验证"},
        {"text": ""},
        {"text": "📅 截止日期：2026-04-07", "bold": True, "highlight": True},
        {"text": "✅ 状态：已完成"},
    ]

    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    for i, item in enumerate(content_items):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()

        p.text = item['text']
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(60, 60, 60)
        p.font.bold = item.get('bold', False)

    return slide

def create_dp5_slide(prs):
    """DP5: 工艺冻结"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "DP5: 工艺冻结"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 0, 0)

    content_items = [
        {"text": "工艺参数确认", "bold": True},
        {"text": "• 起鼓工艺参数优化"},
        {"text": "• 缺陷率控制目标（11.4%）"},
        {"text": "• 4线卷绕质量控制"},
        {"text": ""},
        {"text": "验证完成", "bold": True},
        {"text": "• CBSC现场验证"},
        {"text": "• 无4线卷绕质量问题"},
        {"text": "• BFM4工艺性能改进"},
        {"text": ""},
        {"text": "📅 截止日期：2026-04-07", "bold": True, "highlight": True},
        {"text": "✅ 状态：已完成"},
    ]

    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    for i, item in enumerate(content_items):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()

        p.text = item['text']
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(60, 60, 60)
        p.font.bold = item.get('bold', False)

    return slide

def create_dp6_slide(prs):
    """DP6: 生产发布"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "DP6: 生产发布"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 0, 0)

    content_items = [
        {"text": "量产准备", "bold": True},
        {"text": "• 商业案例最终确认"},
        {"text": "• WWD OEE 验证"},
        {"text": "• 起鼓工艺性能验证"},
        {"text": "• 客户侧工艺性能追踪"},
        {"text": ""},
        {"text": "交付物", "bold": True},
        {"text": "• 项目总结报告"},
        {"text": "• 技术移交文档"},
        {"text": "• 培训材料"},
        {"text": ""},
        {"text": "📅 截止日期：2026-07-30", "bold": True, "highlight": True},
        {"text": "⏰ 状态：进行中（剩余15天）"},
    ]

    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    for i, item in enumerate(content_items):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()

        p.text = item['text']
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(60, 60, 60)
        p.font.bold = item.get('bold', False)

    return slide

def create_dp7_slide(prs):
    """DP7: 移交销售"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "DP7: 移交销售"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 0, 0)

    content_items = [
        {"text": "技术移交", "bold": True},
        {"text": "• 技术文档移交"},
        {"text": "• 培训销售人员"},
        {"text": "• 现场支持计划"},
        {"text": ""},
        {"text": "市场推广", "bold": True},
        {"text": "• 市场推广材料准备"},
        {"text": "• 客户案例收集"},
        {"text": "• 销售支持体系建立"},
        {"text": ""},
        {"text": "📅 截止日期：2026-12-30", "bold": True, "highlight": True},
        {"text": "⏰ 状态：计划中"},
    ]

    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    for i, item in enumerate(content_items):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()

        p.text = item['text']
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(60, 60, 60)
        p.font.bold = item.get('bold', False)

    return slide

def main():
    """主函数"""
    # 创建PPT
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 封面
    create_title_slide(prs, "Bekaert 1+6+12 BFM4-4 Wire Process", "项目汇报", "2026-07-15")

    # 项目概述
    content_items = [
        {"text": "项目名称：1+6+12 BFM4-4 Wire Process", "bold": True},
        {"text": "业务单元：RR"},
        {"text": "项目阶段：Validation（验证阶段）"},
        {"text": "当前状态：✅ Active（活跃）"},
        {"text": ""},
        {"text": "关键人员：", "bold": True},
        {"text": "• 项目发起人：Guanghua Zhang"},
        {"text": "• 项目经理：Xiuhui Li"},
        {"text": ""},
        {"text": "项目目标：", "bold": True},
        {"text": "• 验证 BFM4-4 线材工艺性能"},
        {"text": "• 解决4线卷绕质量问题"},
        {"text": "• 提升起鼓工艺稳定性"},
        {"text": "• 为量产发布做准备"},
    ]
    create_content_slide(prs, "项目概述", content_items)

    # 财务概览
    headers = ["指标", "金额（百万欧元）"]
    rows = [
        ["收入（Revenue）", "€0.00M"],
        ["毛利（Gross Margin）", "€0.28M"],
        ["净现值（NPV）", "€1.8M"],
        ["研发成本（R&D）", "€0.04M"],
        ["资本支出（CAPEX）", "€0.33M"],
    ]
    create_table_slide(prs, "财务概览（X+5预测）", headers, rows)

    # 健康状况
    health_items = [
        {"text": "维度", "bold": True, "indent": True},
        {"text": "进度健康", "indent": True},
        {"text": "✅ On track", "indent": True, "highlight": True},
        {"text": "范围健康", "indent": True},
        {"text": "✅ As defined", "indent": True, "highlight": True},
        {"text": "预算健康", "indent": True},
        {"text": "✅ On target (±10%)", "indent": True, "highlight": True},
    ]
    create_content_slide(prs, "项目健康状况（截至2026-06-30）", health_items)

    # 关键成果
    results_items = [
        {"text": "✅ 成果1：CBSC现场验证完成", "bold": True},
        {"text": "• 设备：BF2+BFM4 SPR2（175吨）"},
        {"text": "• 配置：4台BFM4 + 2台AWC 2.0 WND"},
        {"text": "• 地点：CBSC"},
        {"text": "• 结果："},
        {"text": "  ✓ 无4线卷绕质量问题"},
        {"text": "  ✓ BFM4显示改进的工艺性能"},
        {"text": "  ✓ 缺陷率降至 11.4%"},
        {"text": ""},
        {"text": "✅ 成果2：客户反馈", "bold": True},
        {"text": "• CBSC反馈：BF2/BFM4/BFM6需要更多机器"},
        {"text": "• BFM4目前是产能瓶颈"},
        {"text": "• 市场需求旺盛，需加速推进"},
    ]
    create_content_slide(prs, "关键成果", results_items)

    # 里程碑
    milestone_items = [
        {"text": "节点点", "bold": True, "indent": True},
        {"text": "截止日期", "indent": True, "indent": True},
        {"text": "状态", "indent": True, "indent": True, "indent": True},
        {"text": "项目启动 DP1", "indent": True},
        {"text": "2025-06-30", "indent": True, "indent": True},
        {"text": "✅ 已完成", "indent": True, "indent": True, "indent": True},
        {"text": "概念冻结 DP3", "indent": True},
        {"text": "2026-04-07", "indent": True, "indent": True},
        {"text": "✅ 已完成", "indent": True, "indent": True, "indent": True},
        {"text": "工艺冻结 DP5", "indent": True},
        {"text": "2026-04-07", "indent": True, "indent": True},
        {"text": "✅ 已完成", "indent": True, "indent": True, "indent": True},
        {"text": "生产发布 DP6", "indent": True},
        {"text": "2026-07-30", "indent": True, "indent": True},
        {"text": "⏳ 计划中", "indent": True, "indent": True, "indent": True},
        {"text": "移交销售 DP7", "indent": True},
        {"text": "2026-12-30", "indent": True, "indent": True},
        {"text": "⏳ 计划中", "indent": True, "indent": True, "indent": True},
    ]
    create_content_slide(prs, "项目里程碑", milestone_items)

    # 下一步交付物
    deliverables_items = [
        {"text": "交付物1：客户侧工艺性能追踪", "bold": True},
        {"text": "• 持续收集CBSC现场运行数据"},
        {"text": "• 监控起鼓均值稳定性"},
        {"text": "• 追踪缺陷率变化趋势"},
        {"text": ""},
        {"text": "交付物2：商业案例验证", "bold": True},
        {"text": "• WWD OEE（整体设备效率）验证"},
        {"text": "• 起鼓工艺性能验证"},
        {"text": "• 经济性分析确认"},
        {"text": ""},
        {"text": "交付物3：DP6门评审准备", "bold": True},
        {"text": "• 准备DP6评审材料"},
        {"text": "• 确保所有验证数据完整"},
        {"text": "• 获得管理层批准"},
        {"text": "📅 目标：2026-07-30 完成DP6评审"},
    ]
    create_content_slide(prs, "下一步交付物", deliverables_items)

    # 总结
    summary_items = [
        "项目处于验证阶段，三个维度（进度/范围/预算）均健康",
        "CBSC现场验证取得积极成果，缺陷率降至11.4%",
        "客户反馈BFM4是产能瓶颈，市场需求旺盛",
        "NPV €1.8M，项目经济价值良好",
    ]

    recommendations = [
        "加速推进DP6评审，确保7月30日前完成",
        "重点关注起鼓均值稳定性（当前均值1.64）",
        "继续优化缺陷率，目标降至10%以下",
        "为12月DP7移交销售做好充分准备",
    ]

    create_summary_slide(prs, "总结与建议", summary_items, recommendations)

    # DP1
    create_dp1_slide(prs)

    # DP3
    create_dp3_slide(prs)

    # DP5
    create_dp5_slide(prs)

    # DP6
    create_dp6_slide(prs)

    # DP7
    create_dp7_slide(prs)

    # 保存PPT
    output_path = "/home/YDL/.openclaw/workspace/a_stock_plan/bekaert_project_presentation.pptx"
    prs.save(output_path)

    print(f"✅ PPT已生成：{output_path}")
    print(f"📊 共 {len(prs.slides)} 页")

if __name__ == "__main__":
    main()
